"""Read a PowerPoint file and output text content.
Supports: .pptx, .ppt

Improvements over read-pptx.py:
  - Extracts speaker notes (often contain the most important context)
  - Recursively extracts text from grouped shapes
  - Includes slide layout name for structural context
  - Better .ppt fallback: tries LibreOffice conversion to .pptx then re-reads

Usage: python read-pptx.py <file_path>
"""
import sys
import os
import subprocess

if len(sys.argv) < 2:
    print("Usage: python read-pptx.py <file_path>")
    sys.exit(1)

path = sys.argv[1]
ext = os.path.splitext(path)[1].lower()


def extract_text_from_shape(shape, texts):
    """Recursively extract text from a shape, handling groups."""
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for child in shape.shapes:
                extract_text_from_shape(child, texts)
            return
    except Exception:
        pass

    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            if para.text.strip():
                texts.append(para.text.strip())

    if shape.has_table:
        for row in shape.table.rows:
            texts.append("\t".join(cell.text.strip() for cell in row.cells))

    if shape.has_chart:
        chart = shape.chart
        try:
            title = chart.chart_title.text_frame.text if chart.has_title else "Untitled"
        except Exception:
            title = "Untitled"
        texts.append(f"\n[Chart: {title}]")
        try:
            # Category axis labels
            try:
                cats = [str(c) for c in chart.category_axis.categories]
                if cats:
                    texts.append(f"  Categories: {', '.join(cats)}")
            except Exception:
                pass
            # Series values
            for series in chart.series:
                name = series.name if hasattr(series, 'name') and series.name else "Series"
                vals = [str(v) for v in series.values if v is not None] if series.values else []
                texts.append(f"  {name}: {', '.join(vals)}")
        except Exception as e:
            texts.append(f"  (chart data extraction failed: {e})")


if ext == '.pptx':
    try:
        from pptx import Presentation
        prs = Presentation(path)

        for i, slide in enumerate(prs.slides):
            texts = []

            # Slide layout name gives structural context (e.g. "Title Slide", "Section Header")
            try:
                layout_name = slide.slide_layout.name
                if layout_name:
                    texts.append(f"[Layout: {layout_name}]")
            except Exception:
                pass

            # All shapes -- recursive to catch grouped shapes
            for shape in slide.shapes:
                extract_text_from_shape(shape, texts)

            # Speaker notes -- extract last so they appear after slide content
            notes_text = ""
            try:
                notes_tf = slide.notes_slide.notes_text_frame
                notes_text = notes_tf.text.strip() if notes_tf else ""
            except Exception:
                pass

            if texts or notes_text:
                print(f"--- Slide {i + 1} ---")
                if texts:
                    print("\n".join(texts))
                if notes_text:
                    print(f"\n[Speaker Notes]\n{notes_text}")
                print()

        sys.exit(0)
    except ImportError:
        print("[ERROR] Cannot read PPTX: install python-pptx\n  pip install python-pptx")
        sys.exit(1)
    except Exception as e:
        print(f"[ERROR] Failed to read PPTX: {e}")
        sys.exit(1)

elif ext == '.ppt':
    # Strategy 1: catppt (Linux/Mac tool)
    try:
        result = subprocess.run(['catppt', path], capture_output=True, text=True, timeout=30)
        if result.returncode == 0 and result.stdout.strip():
            print(result.stdout)
            sys.exit(0)
    except Exception:
        pass

    # Strategy 2: LibreOffice -- convert to .pptx then re-read with python-pptx
    try:
        import tempfile
        tmp_dir = tempfile.mkdtemp()
        result = subprocess.run(
            ['libreoffice', '--headless', '--convert-to', 'pptx', '--outdir', tmp_dir, path],
            capture_output=True, text=True, timeout=90
        )
        if result.returncode == 0:
            converted = os.path.join(tmp_dir, os.path.splitext(os.path.basename(path))[0] + '.pptx')
            if os.path.exists(converted):
                result2 = subprocess.run(
                    [sys.executable, __file__, converted],
                    capture_output=True, text=True
                )
                print(result2.stdout)
                sys.exit(result2.returncode)
    except Exception:
        pass

    print("[ERROR] Cannot read .ppt format. Install catppt or LibreOffice for legacy .ppt support.")
    sys.exit(1)

else:
    print(f"[ERROR] Unsupported format: {ext}")
    sys.exit(1)
