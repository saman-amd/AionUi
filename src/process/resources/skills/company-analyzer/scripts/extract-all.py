"""Pre-extract all company documents into plain text files.
Creates a _extracted/ folder with readable .txt versions of every document.

Improvements over extract-all.py:
  - Incremental: skips files where .txt output is already newer than source
  - Parallel processing with ThreadPoolExecutor (--workers N, default 4)
  - .doc files use a proper fallback chain (antiword -> docx2txt) instead of
    python-docx which doesn't support the legacy .doc format
  - PPTX extraction includes speaker notes and chart data (parity with read-pptx.py)
  - XLSX extraction handles merged cells properly
  - PDF extraction prefers pdfplumber for better layout/table handling
  - Progress shows skipped/extracted/failed counts separately
  - --force flag to re-extract everything unconditionally

Usage:
  python extract-all.py [directory] [--force] [--workers N]

Options:
  directory   Directory to scan (default: current working directory)
  --force     Re-extract all files even if output is up-to-date
  --workers   Number of parallel worker threads (default: 4)
"""

import os
import sys
import argparse
from concurrent.futures import ThreadPoolExecutor, as_completed


# ---------------------------------------------------------------------------
# Extractors
# ---------------------------------------------------------------------------

def extract_pdf(filepath):
    """Prefer pdfplumber (best layout/table handling), fall back to pypdf."""
    try:
        import pdfplumber
        output = []
        with pdfplumber.open(filepath) as pdf:
            for i, page in enumerate(pdf.pages):
                tables = page.extract_tables()
                text = page.extract_text(x_tolerance=2, y_tolerance=2)
                if not text and not tables:
                    continue
                output.append(f"--- Page {i + 1} ---")
                if tables:
                    for t_idx, table in enumerate(tables):
                        output.append(f"[Table {t_idx + 1}]")
                        for row in table:
                            cells = [str(c).strip() if c is not None else "" for c in row]
                            output.append("\t".join(cells))
                        output.append("")
                if text and text.strip():
                    output.append(text.strip())
                output.append("")
        return "\n".join(output) if output else "[No text content found in PDF]"
    except ImportError:
        pass
    try:
        from pypdf import PdfReader
        reader = PdfReader(filepath)
        if reader.is_encrypted:
            try:
                reader.decrypt("")
            except Exception:
                return "[ERROR] PDF is password-protected"
        output = []
        for i, page in enumerate(reader.pages):
            text = page.extract_text()
            if text and text.strip():
                output.append(f"--- Page {i + 1} ---")
                output.append(text.strip())
                output.append("")
        return "\n".join(output) if output else "[No text content found in PDF]"
    except ImportError:
        return "[ERROR] Cannot read PDF: install pdfplumber or pypdf"


def _resolve_merged_cells(ws):
    """Fill every merged cell with the top-left value of its range."""
    merged = {}
    for merge_range in ws.merged_cells.ranges:
        top_left = ws.cell(merge_range.min_row, merge_range.min_col).value
        for row in range(merge_range.min_row, merge_range.max_row + 1):
            for col in range(merge_range.min_col, merge_range.max_col + 1):
                merged[(row, col)] = top_left
    return merged


def extract_xlsx(filepath):
    """Extract Excel with merged-cell support and no row cap."""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(filepath, data_only=True)
        output = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            merged = _resolve_merged_cells(ws)
            output.append(f"=== Sheet: {sheet_name} ===")
            for row in ws.iter_rows():
                vals = []
                for cell in row:
                    v = merged.get((cell.row, cell.column), cell.value)
                    vals.append(str(v) if v is not None else "")
                if any(v.strip() for v in vals):
                    output.append("\t".join(vals))
            # Public .charts property (not private ._charts)
            try:
                for ci, chart in enumerate(ws.charts):
                    title = getattr(chart, 'title', None) or f"Chart {ci + 1}"
                    if hasattr(title, 'text'):
                        title = title.text
                    output.append(f"\n[Chart {ci + 1}: {title}]")
            except Exception:
                pass
            output.append("")
        return "\n".join(output)
    except ImportError:
        pass
    try:
        import pandas as pd
        pd.set_option('display.max_colwidth', None)
        xls = pd.ExcelFile(filepath)
        output = []
        for sheet in xls.sheet_names:
            df = pd.read_excel(xls, sheet_name=sheet, header=None)
            output.append(f"=== Sheet: {sheet} ===")
            output.append(df.to_string(index=False, header=False, na_rep=""))
            output.append("")
        return "\n".join(output)
    except ImportError:
        return "[ERROR] Cannot read XLSX: install openpyxl or pandas"


def extract_docx(filepath):
    """Extract DOCX including headers and footers."""
    try:
        from docx import Document
        doc = Document(filepath)
        output = []
        # Headers/footers
        hf_texts = []
        seen = set()
        for section in doc.sections:
            for attr in ('header', 'footer', 'first_page_header', 'first_page_footer'):
                try:
                    hf = getattr(section, attr)
                    if hf and not hf.is_linked_to_previous:
                        for p in hf.paragraphs:
                            t = p.text.strip()
                            if t and t not in seen:
                                hf_texts.append(t)
                                seen.add(t)
                except Exception:
                    pass
        if hf_texts:
            output.append("[Header/Footer]")
            output.extend(hf_texts)
            output.append("")
        # Paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                output.append(para.text)
        # Tables
        for i, table in enumerate(doc.tables):
            output.append(f"\n[Table {i + 1}]")
            for row in table.rows:
                output.append("\t".join(cell.text.strip() for cell in row.cells))
        return "\n".join(output) if output else "[No text content found]"
    except ImportError:
        return "[ERROR] Cannot read DOCX: install python-docx"


def extract_doc(filepath):
    """Extract legacy .doc -- proper fallback chain, no python-docx."""
    import subprocess
    # antiword
    try:
        r = subprocess.run(['antiword', filepath], capture_output=True, text=True, timeout=30)
        if r.returncode == 0 and r.stdout.strip():
            return r.stdout
    except Exception:
        pass
    # docx2txt
    try:
        import docx2txt
        text = docx2txt.process(filepath)
        if text and text.strip():
            return text
    except ImportError:
        pass
    return "[ERROR] Cannot read .doc: install antiword or docx2txt"


def extract_pptx(filepath):
    """Extract PPTX including speaker notes and chart data."""
    try:
        from pptx import Presentation
        prs = Presentation(filepath)
        output = []
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            texts.append(para.text.strip())
                if shape.has_table:
                    for row in shape.table.rows:
                        texts.append("\t".join(cell.text.strip() for cell in row.cells))
                if shape.has_chart:
                    chart = shape.chart
                    try:
                        title = chart.chart_title.text_frame.text if chart.has_title else "Untitled"
                    except Exception:
                        title = "Untitled"
                    texts.append(f"\n[Chart: {title}]")
                    try:
                        for series in chart.series:
                            name = series.name if hasattr(series, 'name') and series.name else "Series"
                            vals = [str(v) for v in series.values if v is not None] if series.values else []
                            texts.append(f"  {name}: {', '.join(vals)}")
                    except Exception:
                        pass
            # Speaker notes
            notes_text = ""
            try:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
            except Exception:
                pass
            if texts or notes_text:
                output.append(f"--- Slide {i + 1} ---")
                output.extend(texts)
                if notes_text:
                    output.append(f"\n[Speaker Notes]\n{notes_text}")
                output.append("")
        return "\n".join(output) if output else "[No text content found]"
    except ImportError:
        return "[ERROR] Cannot read PPTX: install python-pptx"


def extract_txt(filepath):
    with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
        return f.read()


EXTRACTORS = {
    '.xlsx': extract_xlsx,
    '.xlsm': extract_xlsx,
    '.pdf': extract_pdf,
    '.docx': extract_docx,
    '.doc': extract_doc,   # uses proper fallback chain, not python-docx
    '.pptx': extract_pptx,
    '.txt': extract_txt,
    '.csv': extract_txt,
    '.md': extract_txt,
}


# ---------------------------------------------------------------------------
# Incremental logic
# ---------------------------------------------------------------------------

def is_up_to_date(source_path, out_path):
    """True if the output .txt exists and is at least as new as the source file."""
    if not os.path.exists(out_path):
        return False
    return os.path.getmtime(out_path) >= os.path.getmtime(source_path)


# ---------------------------------------------------------------------------
# Worker
# ---------------------------------------------------------------------------

def process_file(task):
    filepath, rel_path, out_path, ext, force = task
    if not force and is_up_to_date(filepath, out_path):
        return ('skipped', rel_path)
    try:
        content = EXTRACTORS[ext](filepath)
        header = f"SOURCE: {rel_path}\nTYPE: {ext}\n{'=' * 60}\n\n"
        os.makedirs(os.path.dirname(out_path), exist_ok=True)
        with open(out_path, 'w', encoding='utf-8') as f:
            f.write(header + content)
        return ('ok', rel_path)
    except Exception as e:
        return ('fail', rel_path, str(e))


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(description='Pre-extract company documents to plain text.')
    parser.add_argument('directory', nargs='?', default=os.getcwd())
    parser.add_argument('--force', action='store_true',
                        help='Re-extract all files even if output is up-to-date')
    parser.add_argument('--workers', type=int, default=4,
                        help='Parallel worker threads (default: 4)')
    args = parser.parse_args()

    root_dir = args.directory
    if not os.path.isdir(root_dir):
        print(f"Error: '{root_dir}' is not a directory")
        sys.exit(1)

    extracted_dir = os.path.join(root_dir, '_extracted')
    os.makedirs(extracted_dir, exist_ok=True)

    # Folders to always skip
    SKIP_DIRS = {'_extracted', '.git', '__pycache__', 'node_modules', '.venv', 'venv'}

    # Collect tasks
    tasks = []
    for dirpath, dirnames, filenames in os.walk(root_dir):
        dirnames[:] = [d for d in dirnames if d not in SKIP_DIRS and not d.startswith('.')]
        for filename in filenames:
            if filename.startswith('.') or filename.startswith('~'):
                continue
            ext = os.path.splitext(filename)[1].lower()
            if ext not in EXTRACTORS:
                continue
            filepath = os.path.join(dirpath, filename)
            rel_path = os.path.relpath(filepath, root_dir)
            out_path = os.path.join(extracted_dir, rel_path + '.txt')
            tasks.append((filepath, rel_path, out_path, ext, args.force))

    print(f"Found {len(tasks)} files. Processing with {args.workers} workers...\n")

    success = failed = skipped = 0
    failed_files = []

    with ThreadPoolExecutor(max_workers=args.workers) as executor:
        futures = {executor.submit(process_file, t): t for t in tasks}
        for future in as_completed(futures):
            result = future.result()
            if result[0] == 'ok':
                print(f"  OK:      {result[1]}")
                success += 1
            elif result[0] == 'skipped':
                print(f"  SKIP:    {result[1]}  (up-to-date)")
                skipped += 1
            else:
                print(f"  FAIL:    {result[1]}  -> {result[2]}")
                failed += 1
                failed_files.append(result[1])

    print(f"\nDone! {success} extracted, {skipped} skipped (up-to-date), {failed} failed.")
    if failed_files:
        print("Failed files:")
        for f in failed_files:
            print(f"  - {f}")
    print(f"Output: {extracted_dir}")

    # Index file
    index_path = os.path.join(extracted_dir, '_INDEX.md')
    with open(index_path, 'w', encoding='utf-8') as f:
        f.write("# Extracted Documents Index\n\n")
        f.write(f"Source: {root_dir}\n")
        f.write(f"Extracted: {success} | Skipped: {skipped} | Failed: {failed}\n\n")
        for dirpath, dirnames, filenames in os.walk(extracted_dir):
            dirnames[:] = [d for d in dirnames if not d.startswith('.')]
            for filename in sorted(filenames):
                if filename == '_INDEX.md':
                    continue
                rel = os.path.relpath(os.path.join(dirpath, filename), extracted_dir)
                f.write(f"- {rel}\n")
    print(f"Index: {index_path}")


if __name__ == '__main__':
    main()
